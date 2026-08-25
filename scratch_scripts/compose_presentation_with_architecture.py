from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


SRC = Path("/mnt/c/rudn/2 курс/VKR/презентация/RUDN_VKR_Dudolin_3.pptx")
DST = Path("/mnt/c/rudn/2 курс/VKR/презентация/RUDN_VKR_Dudolin_4_architecture.pptx")

ROOT = Path(__file__).resolve().parents[1]
ARCH_IMG = ROOT / "reports/figures/summary/presentation_architecture_hybrid_cascade_slide.png"
TRANSFER_IMG = ROOT / "reports/figures/summary/presentation_architecture_transfer_protocol_slide.png"


def cover_content(slide) -> None:
    """Cover the old body area while preserving the slide title and footer."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.22),
        Inches(1.02),
        Inches(9.55),
        Inches(3.66),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rect.line.color.rgb = RGBColor(255, 255, 255)


def add_centered_picture(slide, image_path: Path) -> None:
    slide.shapes.add_picture(
        str(image_path),
        Inches(0.90),
        Inches(1.08),
        width=Inches(8.20),
    )


def replace_text(slide, old: str, new: str) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)


def main() -> None:
    prs = Presentation(str(SRC))

    # Slide numbering is 1-based for the presentation, 0-based in python-pptx.
    slide9 = prs.slides[8]
    replace_text(
        slide9,
        "Гибридная каскадная архитектура с общим CNN-энкодером",
        "Гибридная каскадная архитектура с перенесённым CNN-энкодером",
    )
    cover_content(slide9)
    add_centered_picture(slide9, ARCH_IMG)

    slide10 = prs.slides[9]
    cover_content(slide10)
    add_centered_picture(slide10, TRANSFER_IMG)

    DST.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(DST))
    print(DST)


if __name__ == "__main__":
    main()
